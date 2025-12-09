#!/usr/bin/env python3
"""
Generate CCP4i2 Architecture Evolution presentation.
Run: python3 create_architecture_pptx.py
Output: ccp4i2_architecture_evolution.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, code_block=None):
    """Add a content slide with title, bullets, and optional code."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Bullets
    if bullets:
        bullet_top = Inches(1.2)
        bullet_height = Inches(2.5) if code_block else Inches(5)
        bullet_box = slide.shapes.add_textbox(Inches(0.5), bullet_top, Inches(9), bullet_height)
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.space_after = Pt(12)

    # Code block
    if code_block:
        code_top = Inches(4) if bullets else Inches(1.5)
        code_box = slide.shapes.add_textbox(Inches(0.5), code_top, Inches(9), Inches(2.8))
        tf = code_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = code_block
        p.font.size = Pt(12)
        p.font.name = "Courier New"
        # Add background
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = RGBColor(240, 240, 240)

    return slide


def add_comparison_slide(prs, title, left_title, left_content, right_title, right_content):
    """Add a two-column comparison slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Left column header
    left_header = slide.shapes.add_textbox(Inches(0.3), Inches(1.1), Inches(4.5), Inches(0.5))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(180, 80, 80)

    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.6), Inches(4.5), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = left_content
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(255, 240, 240)

    # Right column header
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.5), Inches(0.5))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(80, 120, 80)

    # Right content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.5), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = right_content
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(240, 255, 240)

    return slide


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Table
    cols = len(headers)
    table_rows = len(rows) + 1  # +1 for header
    table = slide.shapes.add_table(table_rows, cols, Inches(0.3), Inches(1.3), Inches(9.4), Inches(0.5 * table_rows)).table

    # Set column widths
    col_width = Inches(9.4 / cols)
    for i in range(cols):
        table.columns[i].width = col_width

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(70, 130, 180)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 245, 245)

    return slide


def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "CCP4i2: Qt to Django",
        "Modernizing Crystallographic Computing Infrastructure"
    )

    # Slide 2: The Challenge
    add_content_slide(
        prs,
        "The Challenge",
        [
            "Legacy Qt architecture: tightly-coupled GUI + database + business logic",
            "PySide2/Qt bindings increasingly difficult to maintain",
            "Procedural GUI code hard to test, extend, and onboard new developers",
            "10+ years of crystallographic wrappers & pipelines to preserve",
            "KEY GOAL: Reduce long-term support burden while serving existing users"
        ]
    )

    # Slide 3: Architecture Overview
    add_table_slide(
        prs,
        "Architecture Overview",
        ["Layer", "Legacy (Qt)", "Modern (Django)"],
        [
            ["Database", "QtSql (raw SQL)", "Django ORM"],
            ["API", "None (in-process)", "REST API (DRF)"],
            ["Frontend", "Qt Widgets (procedural)", "React (declarative)"],
            ["Events", "Signals / Slots", "asyncio"],
            ["Tooling", "Qt MOC, custom build", "npm, pytest, TypeScript"],
        ]
    )

    # Slide 4: Compatibility Layer
    add_content_slide(
        prs,
        "The Compatibility Layer (baselayer)",
        [
            "Automatic environment detection: Qt or Django mode",
            "Qt API stubs enable existing code to run without PySide2",
            "Gradual migration path - wrappers work unchanged",
            "Key insight: same import works in both worlds"
        ],
        code_block="""# Same wrapper code works in Qt OR Django
from ccp4i2.baselayer import QtCore, Signal, Slot

class MyWrapper(CPluginScript):
    finished = Signal(int)  # Works in both environments

    def processOutputFiles(self):
        self.finished.emit(0)  # Qt: real signal, Django: no-op"""
    )

    # Slide 5: Database Comparison
    add_comparison_slide(
        prs,
        "Database: QtSql vs Django ORM",
        "Before (Imperative SQL)",
        """# Manual query construction
db = QSqlDatabase.addDatabase("QSQLITE")
db.setDatabaseName("project.db")
db.open()

query = QSqlQuery()
query.exec_(\"\"\"
    SELECT * FROM jobs
    WHERE status = 1
\"\"\")

while query.next():
    job_id = query.value(0)
    name = query.value(1)
    # Manual type handling...""",
        "After (Declarative ORM)",
        """# Type-safe, optimized queries
from ccp4x.db import models

# Automatic query optimization
jobs = models.Job.objects.filter(
    status=1
).select_related('project')

for job in jobs:
    print(job.uuid)  # Type-safe
    print(job.project.name)  # No N+1

# Migrations handled automatically
# python manage.py migrate"""
    )

    # Slide 6: Interface Definition
    add_comparison_slide(
        prs,
        "GUI: Procedural Qt vs Declarative React",
        "Procedural Qt (*_gui.py)",
        """# phaser_pipeline_gui.py
class phaser_pipeline_gui(CTaskWidget):

  def drawContents(self):
    self.openFolder(title='Input')
    self.drawReflectionPanel()
    self.openFolder(title='Keywords')

  def drawReflectionPanel(self):
    self.createLine(['subtitle', 'Reflections'])
    self.openSubFrame(frame=True)
    self.createLine(['widget', 'F_SIGF'])
    self.createLine([
      'label','Resolution',
      'widget', 'RESOLUTION_LOW',
      'widget', 'RESOLUTION_HIGH'
    ])
    # Manual signal connection
    self.container.inputData.F_SIGF\\
      .dataChanged.connect(self.update)
    self.closeSubFrame()""",
        "Declarative React (Modern)",
        """// aimless_pipe.tsx
export const AimlessPipeInterface = ({job}) => {
  return (
    <Stack spacing={2}>
      <CCP4i2TaskElement
        job={job}
        itemName="inputData.F_SIGF"
      />
      <CCP4i2TaskElement
        job={job}
        itemName="inputData.RESOLUTION"
        qualifiers={{guiLabel: "Resolution"}}
      />
    </Stack>
  );
};

// Components auto-render from schema
// State managed by React hooks
// REST API for data fetching"""
    )

    # Slide 7: CLI Architecture
    add_content_slide(
        prs,
        "CLI: i2run & i2",
        [
            "i2run: Task execution compatible with legacy scripts",
            "i2: Modern resource-oriented interface",
            "Django management commands for administration",
            "Dual-mode: automatically detects Qt vs Django environment"
        ],
        code_block="""# Task execution (both modes)
i2run phaser_pipeline --XYZIN model.pdb --HKLIN data.mtz

# Modern resource management (Django mode)
i2 projects list                    # List all projects
i2 projects create "New Project"    # Create project
i2 jobs 1 tree                      # Show job tree
i2 run phaser_pipeline [opts]       # i2run compatible
i2 files 42 list                    # List job files
i2 report 42                        # Generate report

# Django management
python manage.py migrate            # Database migrations
python manage.py createproject      # Create project"""
    )

    # Slide 8: Frontend Architecture
    add_content_slide(
        prs,
        "Frontend: React / Electron",
        [
            "React 19 + TypeScript - type-safe, component-based architecture",
            "Moorhen integration for 3D structure visualization",
            "Standard tooling: npm, ESLint, Prettier, Jest",
            "Easier onboarding - widely-known frameworks and patterns",
            "Composable, testable components with clear data flow"
        ],
        code_block="""// TypeScript models mirror Django
interface Job {
  uuid: string;
  status: JobStatus;
  project: Project;
  taskname: string;
}

// React hooks for data fetching
const { job, mutate } = useJob(jobId);
const { data: files } = useJobFiles(jobId);

// Declarative UI composition
<JobView job={job}>
  <TaskInterface />
  <MoorhenViewer structure={outputPdb} />
</JobView>"""
    )

    # Slide 9: Execution Flow
    add_comparison_slide(
        prs,
        "Execution Flow",
        "Qt (Synchronous, Coupled)",
        """┌─────────────────────────────┐
│         Qt GUI              │
│  ┌─────────────────────┐    │
│  │    signal emit      │    │
│  └──────────┬──────────┘    │
│             ▼               │
│  ┌─────────────────────┐    │
│  │   CPluginScript     │    │
│  │   (in event loop)   │    │
│  └──────────┬──────────┘    │
│             ▼               │
│  ┌─────────────────────┐    │
│  │   QtSql direct      │◄───┤
│  └─────────────────────┘    │
│             ▼               │
│  ┌─────────────────────┐    │
│  │   signal → GUI      │    │
│  └─────────────────────┘    │
└─────────────────────────────┘""",
        "Django (Async, Decoupled)",
        """┌──────────┐    ┌──────────┐
│  React   │    │   CLI    │
│ Frontend │    │  (i2)    │
└────┬─────┘    └────┬─────┘
     │               │
     ▼               ▼
┌─────────────────────────┐
│      REST API           │
│   (Django + DRF)        │
└───────────┬─────────────┘
            ▼
┌─────────────────────────┐
│   asyncio task runner   │
│   (subprocess mgmt)     │
└───────────┬─────────────┘
            ▼
┌─────────────────────────┐
│   Django ORM            │
│   (SQLite/PostgreSQL)   │
└─────────────────────────┘"""
    )

    # Slide 10: Updated Core Utilities
    add_table_slide(
        prs,
        "Updated Core Utilities",
        ["Function", "Legacy", "Modern"],
        [
            ["Coordinate introspection", "MMDB", "Gemmi"],
            ["Reflection introspection", "Clipper", "Gemmi"],
            ["MTZ column operations", "cmtzsplit, cmtzjoin, cad", "Gemmi"],
            ["French-Wilson statistics", "ctruncate", "Servalcat"],
        ]
    )

    # Slide 11: Benefits & Path Forward
    add_content_slide(
        prs,
        "Benefits & Path Forward",
        [
            "PRESERVED: All wrappers, pipelines, crystallographic logic",
            "MODERNIZED: Core utilities updated (Gemmi, Servalcat)",
            "MAINTAINABILITY: Standard frameworks with large communities",
            "TESTABILITY: pytest, Jest, REST mocking - easier CI/CD",
            "EXTENSIBILITY: Modular architecture, clear separation of concerns",
            "PATH: Gradual migration via baselayer compatibility layer"
        ]
    )

    # Save
    output_path = "ccp4i2_architecture_evolution.pptx"
    prs.save(output_path)
    print(f"Created: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
